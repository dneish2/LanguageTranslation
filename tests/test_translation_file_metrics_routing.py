import sys
from io import BytesIO
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from TranslationBackend import TranslationBackend
from translation_metrics import MetricsCollector


def _docx_table_bytes() -> bytes:
    doc = Document()
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Table source text"
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(left=1_000_000, top=1_000_000, width=4_000_000, height=1_000_000)
    textbox.text_frame.text = "Slide source text"
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def _pdf_bytes(text: str) -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output.getvalue()


def test_docx_table_translation_uses_file_metrics_kwarg(monkeypatch):
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    backend = TranslationBackend()
    observed_metrics = []

    def fake_translate_text(text, target_language, correlation_id=None, file_metrics=None):
        observed_metrics.append(file_metrics)
        return f"translated:{text}"

    monkeypatch.setattr(backend, "translate_text", fake_translate_text)
    monkeypatch.setattr(backend, "calculate_tokens", lambda _text: 0)

    custom_metrics = MetricsCollector()
    backend.process_docx(
        BytesIO(_docx_table_bytes()),
        target_language="Spanish",
        do_translate=True,
        file_metrics=custom_metrics,
    )

    assert observed_metrics == [custom_metrics]


def test_pptx_shape_translation_uses_file_metrics_kwarg(monkeypatch):
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    backend = TranslationBackend()
    observed_metrics = []

    def fake_translate_text(text, target_language, correlation_id=None, file_metrics=None):
        observed_metrics.append(file_metrics)
        return f"translated:{text}"

    monkeypatch.setattr(backend, "translate_text", fake_translate_text)
    monkeypatch.setattr(backend, "calculate_tokens", lambda _text: 0)

    custom_metrics = MetricsCollector()
    backend.process_pptx(
        BytesIO(_pptx_bytes()),
        target_language="Spanish",
        do_translate=True,
        file_metrics=custom_metrics,
    )

    assert observed_metrics == [custom_metrics]


def test_pdf_block_translation_uses_file_metrics_kwarg(monkeypatch):
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    backend = TranslationBackend()
    observed_metrics = []

    def fake_translate_text(text, target_language, correlation_id=None, file_metrics=None):
        observed_metrics.append(file_metrics)
        return f"translated:{text}"

    monkeypatch.setattr(backend, "translate_text", fake_translate_text)
    monkeypatch.setattr(backend, "calculate_tokens", lambda _text: 0)

    custom_metrics = MetricsCollector()
    backend.process_pdf(
        BytesIO(_pdf_bytes("PDF source text")),
        target_language="Spanish",
        do_translate=True,
        file_metrics=custom_metrics,
    )

    assert observed_metrics == [custom_metrics]
