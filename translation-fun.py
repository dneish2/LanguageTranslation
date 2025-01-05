import os
from io import BytesIO
import signal
import atexit
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import docx
import openai
import tiktoken
from nicegui import ui
from threading import Thread
import re


langs = {
   0: "German",
   1: "Russian",
   2: "Portuguese",
   3: "Spanish",
   4: "French",
   5: "Japanese"
}


# OpenAI translation function
def translate_text(text, client, target_language):
   if not text.strip() or "©" in text:
       return text  # Skip translation for empty or copyright text


   def replace_numbers(match):
       return f"__NUM__{match.group(0)}__NUM__"


   text_with_placeholders = re.sub(r'\d+', replace_numbers, text)


   prompt = f"Translate this to {target_language}: {text_with_placeholders}"
   messages = [
       {
           "role": "system",
           "content": f"You are a helpful assistant that translates any text to {target_language}. Provide only the translated text and nothing else.",
       },
       {"role": "user", "content": prompt},
   ]


   try:
       completion = client.chat.completions.create(
           model="gpt-4o", messages=messages, max_tokens=2000
       )
       translated_response = completion.choices[0].message.content


       if "Translate this" in translated_response or not translated_response.strip():
           return text  # Fallback to original text if translation seems incorrect


       translated_text_with_numbers = re.sub(
           r'__NUM__(\d+)__NUM__', r'\1', translated_response
       )


       return translated_text_with_numbers.strip()
   except Exception as e:
       print(f"Translation error: {e}")
       return text  # Return original text in case of an error


def process_shape(shape, client, target_language):
   if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
       for row in shape.table.rows:
           for cell in row.cells:
               if cell.text_frame:
                   translate_and_fit_text(cell.text_frame, client, target_language)
   elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
       for sub_shape in shape.shapes:
           process_shape(sub_shape, client, target_language)
   elif hasattr(shape, "text_frame") and shape.text_frame:
       translate_and_fit_text(shape.text_frame, client, target_language)


def translate_and_fit_text(text_frame, client, target_language):
   for paragraph in text_frame.paragraphs:
       for run in paragraph.runs:
           original_text = run.text
           translated_text = translate_text(original_text, client, target_language)
           run.text = translated_text
   fit_text(text_frame)


def fit_text(text_frame):
   consistent_font_size = 9  # Set a smaller, consistent font size to prevent overlapping
   for paragraph in text_frame.paragraphs:
       for run in paragraph.runs:
           run.font.size = Pt(consistent_font_size)


def update_text_in_presentation(prs, client, target_language, progress_ui, label_ui):
   total_shapes = sum(len(slide.shapes) for slide in prs.slides)
   shape_count = 0
   for slide in prs.slides:
       for shape in slide.shapes:
           process_shape(shape, client, target_language)
           shape_count += 1
           progress_ui.set_value((shape_count / total_shapes) * 100)
           label_ui.text = f"Processing shape {shape_count}/{total_shapes}"
           ui.update()  # Force UI to update
   return prs


def translate_docx(input_stream, output_stream, client, target_language, progress_ui, label_ui):
   doc = docx.Document(input_stream)
   total_paragraphs = len(doc.paragraphs) + sum(
       len(table.rows) * len(table.columns) for table in doc.tables
   )
   paragraph_count = 0


   for paragraph in doc.paragraphs:
       for run in paragraph.runs:
           if run.text.strip():
               translated_text = translate_text(run.text, client, target_language)
               run.text = translated_text
       paragraph_count += 1
       progress_ui.set_value((paragraph_count / total_paragraphs) * 100)
       label_ui.text = f"Processing paragraph {paragraph_count}/{total_paragraphs}"
       ui.update()


   for table in doc.tables:
       for row in table.rows:
           for cell in row.cells:
               for paragraph in cell.paragraphs:
                   for run in paragraph.runs:
                       if run.text.strip():
                           translated_text = translate_text(run.text, client, target_language)
                           run.text = translated_text
                   paragraph_count += 1
                   progress_ui.set_value((paragraph_count / total_paragraphs) * 100)
                   label_ui.text = f"Processing table cell {paragraph_count}/{total_paragraphs}"
                   ui.update()


   doc.save(output_stream)


# Define containers for UI elements
upload_container = None
progress_container = None
result_container = None
stats_container = None


def reset_ui():
   global upload_container, progress_container, result_container, stats_container
   progress_container.clear()
   result_container.clear()
   stats_container.clear()
   initialize_upload_ui()


def initialize_upload_ui():
   global upload_container
   upload_container.clear()
   with upload_container:
       with ui.card().style(
           "background-color: #f5f5f5; padding: 20px; border-radius: 10px;"
       ).classes("text-center"):
           ui.label("Please upload a document or presentation").style(
               "font-size: 20px; color: #333;"
           )
           ui.upload(on_upload=handle_upload).style(
               "margin-top: 20px; color: #fff; background-color: #ff9800; border: none; border-radius: 5px;"
           )


def handle_translation(input_file, file_extension, lang_select_value, file_name):
   global progress_container, result_container, stats_container
   api_key = os.getenv("OPENAI_API_KEY").encode("ascii", errors="ignore").decode("utf-8")
   if not api_key:
       raise ValueError("The OPENAI_API_KEY environment variable is not set.")
   client = openai.OpenAI(api_key=api_key)


   target_language = langs[lang_select_value]


   progress_container.clear()
   progress_ui = ui.circular_progress(value=0.0, max=100, show_value=True).classes("mx-auto").style("color: #ff9800;")
   label_ui = ui.label(f"Preparing to translate to {target_language}...").classes("text-center").style("color: #333;")


   def translation_thread():
       try:
           output = BytesIO()
           if file_extension == "pptx":
               prs = input_file
               translated_presentation = update_text_in_presentation(
                   prs, client, target_language, progress_ui, label_ui
               )
               translated_presentation.save(output)
           elif file_extension == "docx":
               translate_docx(
                   input_file, output, client, target_language, progress_ui, label_ui
               )
           output.seek(0)  # This call should be on the BytesIO object
           show_result(file_name, target_language, output)


           # Calculate token usage and costs
           token_count = calculate_tokens(input_file, file_extension)
           cost_estimate = calculate_cost(token_count)
           show_stats(token_count, cost_estimate)
       except Exception as e:
           show_error(e)


   Thread(target=translation_thread).start()


def show_result(file_name, target_language, output):
   global result_container
   result_container.clear()
   with result_container:
       with ui.card().style(
           "background-color: #f5f5f5; padding: 20px; border-radius: 10px;"
       ).classes("text-center"):
           ui.label(
               f"The file '{file_name}' has been translated to {target_language}!"
           ).style("font-size: 18px; color: #333;")
           ui.button(
               "Download",
               on_click=lambda: ui.download(output.read(), f"translated_{file_name}"),
           ).style(
               "margin-top: 20px; color: #fff; background-color: #ff9800; border: none; border-radius: 5px;"
           )
           ui.button("Reset", on_click=reset_ui).style(
               "margin-top: 10px; color: #fff; background-color: #ff9800; border: none; border-radius: 5px;"
           )


def show_error(error):
   global result_container
   result_container.clear()
   with result_container:
       ui.label(f"An error occurred: {error}").style(
           "font-size: 18px; color: #e53935;"
       )


def parse_upload(content, file_extension):
   source_stream = BytesIO(content)
   if file_extension == "pptx":
       prs = Presentation(source_stream)
       return prs
   elif file_extension == "docx":
       return source_stream


def handle_upload(e):
   global progress_container, result_container
   try:
       file_extension = e.name.split(".")[-1]
       input_file = parse_upload(e.content.read(), file_extension)
       file_name = e.name
       progress_container.clear()
       result_container.clear()
       with ui.card().style(
           "background-color: #f5f5f5; padding: 20px; border-radius: 10px;"
       ).classes("text-center"):
           ui.label(
               "I've read your document, what language would you like me to translate it to?"
           ).style("font-size: 18px; color: #333;")
           lang_select = (
               ui.select(langs, value=1)
               .classes("mx-auto")
               .style(
                   "margin-top: 20px; width: 200px; color: #fff; background-color: #ff9800; border: none; border-radius: 5px; text-align: center;"
               )
           )
           ui.button(
               "Translate!",
               on_click=lambda: handle_translation(
                   input_file, file_extension, lang_select.value, file_name
               ),
           ).style(
               "margin-top: 20px; color: #fff; background-color: #ff9800; border: none; border-radius: 5px;"
           )
   except Exception as err:
       with ui.card().style(
           "background-color: #f5f5f5; padding: 20px; border-radius: 10px;"
       ).classes("text-center"):
           ui.label(
               f"I couldn't quite read this document ({err}), could you please upload another?"
           ).style("font-size: 18px; color: #333;")
           ui.upload(on_upload=handle_upload).style(
               "margin-top: 20px; color: #fff; background-color: #ff9800; border: none; border-radius: 5px;"
           )


def calculate_tokens(input_file, file_extension):
   try:
       encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
   except KeyError:
       encoding = tiktoken.get_encoding("cl100k_base")  # Fallback encoding
  
   total_text = ""
  
   if file_extension == "pptx":
       for slide in input_file.slides:
           for shape in slide.shapes:
               if hasattr(shape, "text") and shape.text:
                   total_text += shape.text
   elif file_extension == "docx":
       doc = docx.Document(input_file)
       for paragraph in doc.paragraphs:
           total_text += paragraph.text
       for table in doc.tables:
           for row in table.rows:
               for cell in row.cells:
                   total_text += cell.text


   num_tokens = len(encoding.encode(total_text))
   return num_tokens


def calculate_cost(token_count):
   cost_per_1000_tokens = 0.002
   return (token_count / 1000) * cost_per_1000_tokens


def show_stats(token_count, cost_estimate):
   global stats_container
   stats_container.clear()
   with stats_container:
       with ui.card().style(
           "background-color: #f5f5f5; padding: 20px; border-radius: 10px;"
       ).classes("text-center"):
           ui.label(f"Token usage: {token_count} tokens").style(
               "font-size: 18px; color: #333;"
           )
           ui.label(f"Estimated cost: ${cost_estimate:.4f}").style(
               "font-size: 18px; color: #333;"
           )


@ui.page('/')
def main_page():
   global upload_container, progress_container, result_container, stats_container
   with ui.row().classes("justify-center") as main_container:
       with ui.column().classes("justify-center items-center") as upload_container:
           initialize_upload_ui()
       with ui.column().classes("justify-center items-center") as progress_container:
           pass  # Container for progress UI elements
       with ui.column().classes("justify-center items-center") as result_container:
           pass  # Container for result UI elements
       with ui.column().classes("justify-center items-center") as stats_container:
           pass  # Container for statistics UI elements


def on_shutdown():
   print("Shutting down...")
   try:
       ui.run_javascript("window.close();")  # Close the browser tab
   except RuntimeError as e:
       print(f"Error in on_shutdown: {e}")
   os._exit(0)


if __name__ in {"__main__", "__mp_main__"}:
   atexit.register(on_shutdown)
   signal.signal(signal.SIGINT, lambda s, f: on_shutdown())
   signal.signal(signal.SIGTERM, lambda s, f: on_shutdown())
   ui.run(port=3030)
